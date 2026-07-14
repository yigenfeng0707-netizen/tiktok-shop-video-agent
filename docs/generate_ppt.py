# -*- coding: utf-8 -*-
"""
路演 PPT 生成脚本（图文并茂版）
生成 5 页路演 PPT，嵌入 docs/charts/ 下 13 张图表。
运行：python docs/generate_ppt.py（在 "项目" 目录下运行）
输出：docs/路演PPT.pptx
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# ============ 主题色 ============
COLOR_DARK_BLUE = RGBColor(0x1F, 0x38, 0x64)   # 深蓝 主色
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)       # 白色
COLOR_ORANGE = RGBColor(0xED, 0x7D, 0x31)      # 橙色 强调
COLOR_LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)  # 浅灰 背景块
COLOR_DARK_TEXT = RGBColor(0x33, 0x33, 0x33)   # 深色正文
COLOR_FOOTER_GRAY = RGBColor(0xE7, 0xEA, 0xF0) # 页脚浅灰

# ============ 幻灯片尺寸（16:9） ============
SLIDE_WIDTH = Emu(12192000)
SLIDE_HEIGHT = Emu(6858000)

# ============ 布局常量 ============
TITLE_BAR_H = Emu(760000)
ACCENT_H = Emu(60000)
CONTENT_TOP = Emu(950000)
FOOTER_TOP = Emu(6510000)
FOOTER_H = Emu(348000)

FONT_NAME = "微软雅黑"
CHARTS_DIR = "docs/charts"


# ============ 通用辅助函数 ============
def set_run(run, text, size=18, bold=False, color=COLOR_DARK_TEXT, font=FONT_NAME):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title_bar(slide, title_text, page_num):
    """顶部标题栏：深蓝背景白字 + 橙色装饰条 + 右侧页码"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, TITLE_BAR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_DARK_BLUE
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.margin_left = Emu(400000)
    tf.margin_top = Emu(80000)
    tf.margin_bottom = Emu(80000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run(run, title_text, size=20, bold=True, color=COLOR_WHITE)

    # 橙色装饰条
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, TITLE_BAR_H, SLIDE_WIDTH, ACCENT_H)
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_ORANGE
    accent.line.fill.background()
    accent.shadow.inherit = False

    # 页码（标题栏右侧）
    page_box = slide.shapes.add_textbox(SLIDE_WIDTH - Emu(1300000), 0, Emu(1200000), TITLE_BAR_H)
    ptf = page_box.text_frame
    ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ptf.margin_top = 0
    ptf.margin_bottom = 0
    pp = ptf.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    prun = pp.add_run()
    set_run(prun, f"{page_num} / 5", size=12, bold=True, color=COLOR_WHITE)


def add_footer(slide):
    """底部页脚：浅灰背景，左项目名 / 右团队信息"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, FOOTER_TOP, SLIDE_WIDTH, FOOTER_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_FOOTER_GRAY
    bar.line.fill.background()
    bar.shadow.inherit = False

    lb = slide.shapes.add_textbox(Emu(300000), FOOTER_TOP, Emu(8000000), FOOTER_H)
    ltf = lb.text_frame
    ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ltf.margin_top = 0
    ltf.margin_bottom = 0
    lp = ltf.paragraphs[0]
    lp.alignment = PP_ALIGN.LEFT
    lr = lp.add_run()
    set_run(lr, "TikTok Shop 跨境短视频 AI 本地化 Agent", size=9, color=COLOR_DARK_BLUE)

    rb = slide.shapes.add_textbox(SLIDE_WIDTH - Emu(4300000), FOOTER_TOP, Emu(4000000), FOOTER_H)
    rtf = rb.text_frame
    rtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rtf.margin_top = 0
    rtf.margin_bottom = 0
    rp = rtf.paragraphs[0]
    rp.alignment = PP_ALIGN.RIGHT
    rr = rp.add_run()
    set_run(rr, "Solo · 冯亦根", size=9, color=COLOR_DARK_BLUE)


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=COLOR_DARK_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(80000)
    tf.margin_right = Emu(80000)
    tf.margin_top = Emu(40000)
    tf.margin_bottom = Emu(40000)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color)
    return box


def add_multiline(slide, left, top, width, height, lines, size=12,
                  color=COLOR_DARK_TEXT, align=PP_ALIGN.LEFT, bold_first=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(80000)
    tf.margin_right = Emu(80000)
    tf.margin_top = Emu(40000)
    tf.margin_bottom = Emu(40000)
    for j, line in enumerate(lines):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        set_run(run, line, size=size, color=color, bold=(bold_first and j == 0))
    return box


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    if line_color is not None:
        rect.line.color.rgb = line_color
        rect.line.width = Pt(1)
    else:
        rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    if line_color is not None:
        rect.line.color.rgb = line_color
        rect.line.width = Pt(1.5)
    else:
        rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def shape_text(shape, text, size, bold, color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(60000)
    tf.margin_right = Emu(60000)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color)


def add_image_fit(slide, img_name, box_left, box_top, box_w, box_h):
    """在指定 box 内按比例缩放并居中嵌入图片，保持宽高比不变形。"""
    img_path = f"{CHARTS_DIR}/{img_name}"
    with Image.open(img_path) as im:
        iw, ih = im.size
    ratio = min(box_w / iw, box_h / ih)
    w = int(iw * ratio)
    h = int(ih * ratio)
    left = box_left + (box_w - w) // 2
    top = box_top + (box_h - h) // 2
    return slide.shapes.add_picture(img_path, left, top, width=w, height=h)


def add_stat_cell(slide, left, top, w, h, number, label, bg=COLOR_DARK_BLUE):
    """统计单元格：大数字 + 小标签"""
    add_rounded_rect(slide, left, top, w, h, bg)
    add_textbox(slide, left, top + Emu(40000), w, Emu(230000),
                number, size=20, bold=True, color=COLOR_ORANGE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, left, top + Emu(250000), w, Emu(h - 260000),
                label, size=11, color=COLOR_WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============ 创建演示文稿 ============
def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


# ============ P1: 问题场景 ============
def add_slide_problem(prs):
    """P1: TikTok Shop 跨境卖家的视频生产困境"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "TikTok Shop 跨境卖家的视频生产困境", 1)
    add_footer(slide)

    # 左侧：4 大痛点对比图（大图）
    add_image_fit(slide, "p1_pain_points.png",
                  Emu(300000), Emu(950000), Emu(5900000), Emu(5000000))

    # 右上：市场规模图
    add_image_fit(slide, "p1_market_size.png",
                  Emu(6350000), Emu(950000), Emu(5550000), Emu(2350000))

    # 右下：用户画像
    add_image_fit(slide, "p1_user_persona.png",
                  Emu(6350000), Emu(3400000), Emu(5550000), Emu(2550000))

    # 底部：关键数据条
    data_top = Emu(6050000)
    data_h = Emu(420000)
    cell_w = Emu(3900000)
    gap = Emu(150000)
    start_left = (SLIDE_WIDTH - 3 * cell_w - 2 * gap) // 2
    stats = [("50-80 万", "跨境卖家规模"),
             ("250-400 亿", "视频制作年支出（元）"),
             ("98%", "成本降低")]
    for i, (num, lab) in enumerate(stats):
        add_stat_cell(slide, start_left + i * (cell_w + gap), data_top, cell_w, data_h, num, lab)


# ============ P2: 解决方案 ============
def add_slide_solution(prs):
    """P2: AI 本地化短视频生产流水线"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "AI 本地化短视频生产流水线", 2)
    add_footer(slide)

    # 顶部：4 步流水线架构图（横跨全宽）
    add_image_fit(slide, "p2_pipeline.png",
                  Emu(300000), Emu(950000), Emu(11590000), Emu(2900000))

    # 左下：传统 vs AI 雷达图
    add_image_fit(slide, "p2_comparison.png",
                  Emu(300000), Emu(3950000), Emu(4200000), Emu(2100000))

    # 右下：3 大创新点表格
    add_image_fit(slide, "p2_innovation.png",
                  Emu(4650000), Emu(3950000), Emu(7250000), Emu(2100000))

    # 底部：创新点摘要
    bar = add_rounded_rect(slide, Emu(300000), Emu(6120000), Emu(11590000), Emu(380000),
                           COLOR_DARK_BLUE)
    shape_text(bar, "多国批量本地化  ·  文化禁忌自动规避  ·  数据反哺闭环迭代",
               size=14, bold=True, color=COLOR_WHITE)


# ============ P3: Demo 路径 ============
def add_slide_demo(prs):
    """P3: 现场演示"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "现场 Demo 演示", 3)
    add_footer(slide)

    # 左侧：演示流程图（竖向）
    add_image_fit(slide, "p3_demo_flow.png",
                  Emu(350000), Emu(950000), Emu(4300000), Emu(5000000))

    # 右侧标题
    add_textbox(slide, Emu(4850000), Emu(1000000), Emu(7000000), Emu(360000),
                "5 国视频产出矩阵", size=15, bold=True, color=COLOR_DARK_BLUE)

    # 右侧：5 国产出矩阵
    add_image_fit(slide, "p3_countries.png",
                  Emu(4850000), Emu(1400000), Emu(7000000), Emu(3700000))

    # 底部：MCP 链接 + 响应时间 + 工具数
    bar = add_rounded_rect(slide, Emu(350000), Emu(6000000), Emu(11490000), Emu(470000),
                           COLOR_ORANGE)
    shape_text(bar,
               "MCP: open.hirebox.cn/mcp/65_my-agent_4Dsd7ZgAud4"
               "   ·   响应 0.589s   ·   6 工具全部可用",
               size=12, bold=True, color=COLOR_WHITE)


# ============ P4: 技术架构 ============
def add_slide_tech(prs):
    """P4: 技术架构与 AI 使用披露"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "技术架构与 AI 使用披露", 4)
    add_footer(slide)

    # 左侧：技术栈架构图
    add_image_fit(slide, "p4_tech_stack.png",
                  Emu(300000), Emu(950000), Emu(6200000), Emu(5100000))

    # 右上：AI 披露饼图
    add_image_fit(slide, "p4_ai_disclosure.png",
                  Emu(6700000), Emu(950000), Emu(5200000), Emu(2400000))

    # 右下：技术栈文字列表
    list_left = Emu(6700000)
    list_top = Emu(3500000)
    list_w = Emu(5200000)
    list_h = Emu(2550000)
    add_rounded_rect(slide, list_left, list_top, list_w, list_h, COLOR_LIGHT_GRAY)
    add_rect(slide, list_left, list_top, list_w, Emu(70000), COLOR_ORANGE)
    add_textbox(slide, list_left + Emu(180000), list_top + Emu(120000),
                list_w - Emu(280000), Emu(380000),
                "核心技术栈", size=14, bold=True, color=COLOR_DARK_BLUE)
    add_multiline(slide, list_left + Emu(180000), list_top + Emu(560000),
                  list_w - Emu(280000), Emu(1900000),
                  ["• MiXer AI — Agent 主编排",
                   "• MCP — 模型上下文协议",
                   "• GPT-4 — 脚本生成 / 数据复盘",
                   "• HeyGen — 多语种数字人合成",
                   "• ffmpeg — 批量剪辑 / 字幕 BGM",
                   "• python-pptx — 路演物料生成",
                   "• matplotlib — 数据可视化"],
                  size=12, color=COLOR_DARK_TEXT)

    # 底部：AI 技术评分
    bar = add_rounded_rect(slide, Emu(300000), Emu(6120000), Emu(11590000), Emu(380000),
                           COLOR_DARK_BLUE)
    shape_text(bar, "AI 技术评分：100 / 100   ·   A 级",
               size=15, bold=True, color=COLOR_ORANGE)


# ============ P5: 商业化 ============
def add_slide_business(prs):
    """P5: 商业化路径与收入预测"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "商业化路径与收入预测", 5)
    add_footer(slide)

    # 左上：商业模式对比
    add_image_fit(slide, "p5_business_models.png",
                  Emu(300000), Emu(950000), Emu(5800000), Emu(2400000))

    # 右上：收入预测
    add_image_fit(slide, "p5_revenue_forecast.png",
                  Emu(6300000), Emu(950000), Emu(5600000), Emu(2400000))

    # 底部：迭代路径时间线（横跨全宽）
    add_image_fit(slide, "p5_roadmap.png",
                  Emu(300000), Emu(3500000), Emu(11590000), Emu(2550000))

    # 底部数据条
    data_top = Emu(6120000)
    data_h = Emu(380000)
    cell_w = Emu(3900000)
    gap = Emu(150000)
    start_left = (SLIDE_WIDTH - 3 * cell_w - 2 * gap) // 2
    stats = [("50-100 万", "首年收入（元）"),
             ("3 种", "商业模式"),
             ("Q3", "MVP 上线")]
    for i, (num, lab) in enumerate(stats):
        add_stat_cell(slide, start_left + i * (cell_w + gap), data_top, cell_w, data_h, num, lab)


# ============ 主函数 ============
def main():
    prs = create_presentation()
    add_slide_problem(prs)
    add_slide_solution(prs)
    add_slide_demo(prs)
    add_slide_tech(prs)
    add_slide_business(prs)

    output_path = "docs/路演PPT.pptx"
    prs.save(output_path)
    print(f"PPT 已生成：{output_path}")
    print(f"共 {len(prs.slides)} 页")
    total_pics = 0
    for i, slide in enumerate(prs.slides, 1):
        pics = sum(1 for s in slide.shapes if s.shape_type == 13)
        total_pics += pics
        print(f"  第{i}页: {len(slide.shapes)} 个元素, {pics} 张图片")
    print(f"图片总数: {total_pics}")


if __name__ == "__main__":
    main()
